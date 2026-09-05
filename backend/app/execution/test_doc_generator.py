import io
import zipfile
from datetime import datetime, timezone
from typing import Optional, List, Dict, Any

from app.schemas.doc_generator_schemas import (
    DocumentContentModel,
    DocSection,
    PptxSlide,
    PptxSlideCard,
    DocTable,
    DocCallout
)

# ---------------------------------------------------------------------------
# COLOR THEMES
# ---------------------------------------------------------------------------
THEMES = {
    "corporate_blue": {
        "primary_hex": "1E3A8A",      # Deep Navy
        "secondary_hex": "4F46E5",    # Indigo
        "accent_hex": "0284C7",       # Sky Blue
        "dark_hex": "0F172A",         # Slate 900
        "light_bg_hex": "F8FAFC",     # Slate 50
        "border_hex": "CBD5E1",       # Slate 300
        "success_hex": "059669",      # Emerald
        "warning_hex": "D97706",      # Amber
        "critical_hex": "DC2626",     # Rose
    },
    "slate_dark": {
        "primary_hex": "0F172A",
        "secondary_hex": "334155",
        "accent_hex": "6366F1",
        "dark_hex": "020617",
        "light_bg_hex": "F1F5F9",
        "border_hex": "94A3B8",
        "success_hex": "10B981",
        "warning_hex": "F59E0B",
        "critical_hex": "EF4444",
    },
    "emerald_teal": {
        "primary_hex": "065F46",
        "secondary_hex": "0D9488",
        "accent_hex": "14B8A6",
        "dark_hex": "064E3B",
        "light_bg_hex": "F0FDF4",
        "border_hex": "A7F3D0",
        "success_hex": "059669",
        "warning_hex": "D97706",
        "critical_hex": "DC2626",
    }
}


def hex_to_rgb(hex_str: str):
    hex_str = hex_str.lstrip('#')
    return tuple(int(hex_str[i:i+2], 16) for i in (0, 2, 4))


# ===========================================================================
# 1. DOCX BUILDER (Word Document Engine)
# ===========================================================================
class DocxBuilder:
    @classmethod
    def build(cls, content: DocumentContentModel, theme_name: str = "corporate_blue") -> io.BytesIO:
        import docx
        from docx.shared import Inches, Pt, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.enum.table import WD_TABLE_ALIGNMENT, WD_ALIGN_VERTICAL
        from docx.oxml import OxmlElement, parse_xml
        from docx.oxml.ns import nsdecls, qn

        theme = THEMES.get(theme_name, THEMES["corporate_blue"])
        p_rgb = hex_to_rgb(theme["primary_hex"])
        s_rgb = hex_to_rgb(theme["secondary_hex"])
        d_rgb = hex_to_rgb(theme["dark_hex"])

        doc = docx.Document()

        # Set 1-inch standard margins
        sections = doc.sections
        for s in sections:
            s.top_margin = Inches(1.0)
            s.bottom_margin = Inches(1.0)
            s.left_margin = Inches(1.0)
            s.right_margin = Inches(1.0)
            s.different_first_page_header_footer = True

        # Header / Footer for inner pages
        body_section = doc.sections[0]
        footer = body_section.footer
        f_p = footer.paragraphs[0]
        f_p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        f_run = f_p.add_run(f"EVAL AI Platform  |  {content.meta.classification or 'QA Internal'}")
        f_run.font.size = Pt(8.5)
        f_run.font.color.rgb = RGBColor(148, 163, 184)

        # -------------------------------------------------------------
        # Cover Banner / Header Block
        # -------------------------------------------------------------
        title_p = doc.add_paragraph()
        title_p.paragraph_format.space_before = Pt(10)
        title_p.paragraph_format.space_after = Pt(4)
        t_run = title_p.add_run(content.meta.title or "Test Engineering Specification")
        t_run.font.name = "Arial"
        t_run.font.size = Pt(24)
        t_run.font.bold = True
        t_run.font.color.rgb = RGBColor(*p_rgb)

        if content.meta.subtitle:
            sub_p = doc.add_paragraph()
            sub_p.paragraph_format.space_after = Pt(12)
            sub_run = sub_p.add_run(content.meta.subtitle)
            sub_run.font.name = "Arial"
            sub_run.font.size = Pt(13)
            sub_run.font.color.rgb = RGBColor(*s_rgb)

        # Metadata Table Block
        meta_table = doc.add_table(rows=2, cols=2)
        meta_table.alignment = WD_TABLE_ALIGNMENT.CENTER
        meta_table.autofit = False

        date_val = content.meta.date_str or datetime.now(timezone.utc).strftime("%B %d, %Y")
        meta_items = [
            ("Author / Engine", content.meta.author or "EVAL AI Agent"),
            ("Organization", content.meta.organization or "Quality Engineering"),
            ("Generated Date", date_val),
            ("Classification", content.meta.classification or "Confidential"),
        ]

        for i, (k, v) in enumerate(meta_items):
            r_idx = i // 2
            c_idx = i % 2
            cell = meta_table.cell(r_idx, c_idx)
            # Background shading
            shd = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{theme["light_bg_hex"]}"/>')
            cell._tc.get_or_add_tcPr().append(shd)
            p = cell.paragraphs[0]
            p.paragraph_format.space_before = Pt(3)
            p.paragraph_format.space_after = Pt(3)
            k_run = p.add_run(f"{k}: ")
            k_run.font.bold = True
            k_run.font.size = Pt(9.5)
            k_run.font.color.rgb = RGBColor(*d_rgb)
            v_run = p.add_run(v)
            v_run.font.size = Pt(9.5)
            v_run.font.color.rgb = RGBColor(100, 116, 139)

        doc.add_paragraph().paragraph_format.space_after = Pt(8)

        # -------------------------------------------------------------
        # Executive Summary Callout
        # -------------------------------------------------------------
        if content.executive_summary:
            exec_table = doc.add_table(rows=1, cols=1)
            exec_table.alignment = WD_TABLE_ALIGNMENT.CENTER
            exec_cell = exec_table.cell(0, 0)
            shd = parse_xml(f'<w:shd {nsdecls("w")} w:fill="F1F5F9"/>')
            exec_cell._tc.get_or_add_tcPr().append(shd)
            
            p = exec_cell.paragraphs[0]
            p.paragraph_format.space_before = Pt(8)
            p.paragraph_format.space_after = Pt(4)
            h_run = p.add_run("EXECUTIVE SUMMARY\n")
            h_run.font.bold = True
            h_run.font.size = Pt(11)
            h_run.font.color.rgb = RGBColor(*p_rgb)

            b_run = p.add_run(content.executive_summary)
            b_run.font.size = Pt(10)
            b_run.font.color.rgb = RGBColor(*d_rgb)
            doc.add_paragraph().paragraph_format.space_after = Pt(8)

        # -------------------------------------------------------------
        # Sections Hierarchy
        # -------------------------------------------------------------
        for sec_idx, sec in enumerate(content.sections, start=1):
            h_p = doc.add_paragraph()
            h_p.paragraph_format.space_before = Pt(14)
            h_p.paragraph_format.space_after = Pt(4)
            
            h_run = h_p.add_run(f"{sec_idx}. {sec.heading}")
            h_run.font.name = "Arial"
            h_run.font.bold = True
            h_run.font.size = Pt(14 if sec.level == 1 else 12)
            h_run.font.color.rgb = RGBColor(*p_rgb)

            if sec.summary:
                sum_p = doc.add_paragraph()
                sum_p.paragraph_format.space_after = Pt(4)
                s_run = sum_p.add_run(sec.summary)
                s_run.font.italic = True
                s_run.font.size = Pt(10)
                s_run.font.color.rgb = RGBColor(100, 116, 139)

            for para in sec.paragraphs:
                p = doc.add_paragraph()
                p.paragraph_format.space_after = Pt(4)
                p_run = p.add_run(para)
                p_run.font.size = Pt(10)
                p_run.font.color.rgb = RGBColor(*d_rgb)

            for bullet in sec.bullet_points:
                bp = doc.add_paragraph(style='List Bullet')
                bp.paragraph_format.space_after = Pt(2)
                b_run = bp.add_run(bullet)
                b_run.font.size = Pt(10)
                b_run.font.color.rgb = RGBColor(*d_rgb)

            # Callout Boxes in Section
            for callout in (sec.callouts or []):
                c_tbl = doc.add_table(rows=1, cols=1)
                c_tbl.alignment = WD_TABLE_ALIGNMENT.CENTER
                c_cell = c_tbl.cell(0, 0)
                bg_color = "FEF3C7" if callout.type == "warning" else "EFF6FF" if callout.type == "info" else "F0FDF4"
                shd = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{bg_color}"/>')
                c_cell._tc.get_or_add_tcPr().append(shd)

                cp = c_cell.paragraphs[0]
                cp.paragraph_format.space_before = Pt(4)
                cp.paragraph_format.space_after = Pt(4)
                if callout.title:
                    ct_run = cp.add_run(f"[{callout.type.upper()}] {callout.title}\n")
                    ct_run.font.bold = True
                    ct_run.font.size = Pt(9.5)
                cb_run = cp.add_run(callout.content)
                cb_run.font.size = Pt(9.5)
                doc.add_paragraph().paragraph_format.space_after = Pt(4)

            # Tables in Section
            for tbl in (sec.tables or []):
                if tbl.caption:
                    cap_p = doc.add_paragraph()
                    cap_p.paragraph_format.space_before = Pt(6)
                    cap_p.paragraph_format.space_after = Pt(2)
                    c_run = cap_p.add_run(f"Table: {tbl.caption}")
                    c_run.font.bold = True
                    c_run.font.size = Pt(9.5)

                num_rows = len(tbl.rows) + (1 if tbl.headers else 0)
                num_cols = max(len(tbl.headers), max((len(r) for r in tbl.rows), default=1))
                
                if num_rows > 0 and num_cols > 0:
                    t = doc.add_table(rows=num_rows, cols=num_cols)
                    t.alignment = WD_TABLE_ALIGNMENT.CENTER
                    
                    row_offset = 0
                    if tbl.headers:
                        for col_idx, h_text in enumerate(tbl.headers):
                            cell = t.cell(0, col_idx)
                            shd = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{theme["primary_hex"]}"/>')
                            cell._tc.get_or_add_tcPr().append(shd)
                            hp = cell.paragraphs[0]
                            hp.alignment = WD_ALIGN_PARAGRAPH.CENTER
                            hrun = hp.add_run(str(h_text))
                            hrun.font.bold = True
                            hrun.font.size = Pt(9.5)
                            hrun.font.color.rgb = RGBColor(255, 255, 255)
                        row_offset = 1

                    for r_idx, r_data in enumerate(tbl.rows):
                        zebra_bg = "F8FAFC" if r_idx % 2 == 1 else "FFFFFF"
                        for c_idx, val in enumerate(r_data):
                            if c_idx < num_cols:
                                cell = t.cell(r_idx + row_offset, c_idx)
                                shd = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{zebra_bg}"/>')
                                cell._tc.get_or_add_tcPr().append(shd)
                                dp = cell.paragraphs[0]
                                drun = dp.add_run(str(val))
                                drun.font.size = Pt(9)
                                drun.font.color.rgb = RGBColor(*d_rgb)

                    doc.add_paragraph().paragraph_format.space_after = Pt(6)

        stream = io.BytesIO()
        doc.save(stream)
        stream.seek(0)
        return stream


# ===========================================================================
# 2. PDF BUILDER (ReportLab Engine with NumberedCanvas)
# ===========================================================================
class NumberedCanvas:
    """Two-pass canvas to dynamically compute and write 'Page X of Y' footers."""
    pass  # Dynamically generated inside build


class PdfBuilder:
    @classmethod
    def build(cls, content: DocumentContentModel, theme_name: str = "corporate_blue") -> io.BytesIO:
        from reportlab.lib import colors
        from reportlab.lib.pagesizes import letter
        from reportlab.lib.units import inch
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.platypus import (
            SimpleDocTemplate,
            Paragraph,
            Spacer,
            Table,
            TableStyle,
            KeepTogether,
            HRFlowable
        )
        from reportlab.pdfgen import canvas

        theme = THEMES.get(theme_name, THEMES["corporate_blue"])
        primary_color = colors.HexColor(f"#{theme['primary_hex']}")
        secondary_color = colors.HexColor(f"#{theme['secondary_hex']}")
        dark_color = colors.HexColor(f"#{theme['dark_hex']}")
        light_bg = colors.HexColor(f"#{theme['light_bg_hex']}")
        border_color = colors.HexColor(f"#{theme['border_hex']}")

        class PageNumCanvas(canvas.Canvas):
            def __init__(self, *args, **kwargs):
                super().__init__(*args, **kwargs)
                self.pages = []

            def showPage(self):
                self.pages.append(dict(self.__dict__))
                self._startPage()

            def save(self):
                page_count = len(self.pages)
                for page in self.pages:
                    self.__dict__.update(page)
                    self.draw_page_decorations(page_count)
                    super().showPage()
                super().save()

            def draw_page_decorations(self, page_count):
                if self._pageNumber > 1:
                    # Running Header
                    self.saveState()
                    self.setFont("Helvetica", 8)
                    self.setFillColor(colors.HexColor("#64748B"))
                    self.drawString(54, 11 * inch - 36, (content.meta.title or "Test Engineering Spec")[:60])
                    self.setStrokeColor(colors.HexColor("#CBD5E1"))
                    self.setLineWidth(0.5)
                    self.line(54, 11 * inch - 42, 8.5 * inch - 54, 11 * inch - 42)
                    self.restoreState()

                # Running Footer
                self.saveState()
                self.setFont("Helvetica", 8)
                self.setFillColor(colors.HexColor("#64748B"))
                self.drawString(54, 36, f"EVAL AI Platform  |  {content.meta.classification or 'QA Internal'}")
                page_str = f"Page {self._pageNumber} of {page_count}"
                self.drawRightString(8.5 * inch - 54, 36, page_str)
                self.setStrokeColor(colors.HexColor("#CBD5E1"))
                self.setLineWidth(0.5)
                self.line(54, 48, 8.5 * inch - 54, 48)
                self.restoreState()

        stream = io.BytesIO()
        doc = SimpleDocTemplate(
            stream,
            pagesize=letter,
            leftMargin=54,
            rightMargin=54,
            topMargin=54,
            bottomMargin=54
        )

        styles = getSampleStyleSheet()
        title_style = ParagraphStyle(
            'DocTitle',
            parent=styles['Heading1'],
            fontName='Helvetica-Bold',
            fontSize=22,
            leading=26,
            textColor=primary_color,
            spaceAfter=4
        )
        sub_style = ParagraphStyle(
            'DocSubtitle',
            parent=styles['Normal'],
            fontName='Helvetica',
            fontSize=12,
            leading=16,
            textColor=secondary_color,
            spaceAfter=12
        )
        h1_style = ParagraphStyle(
            'SectionH1',
            parent=styles['Heading2'],
            fontName='Helvetica-Bold',
            fontSize=13,
            leading=17,
            textColor=primary_color,
            spaceBefore=12,
            spaceAfter=4
        )
        body_style = ParagraphStyle(
            'DocBody',
            parent=styles['Normal'],
            fontName='Helvetica',
            fontSize=9.5,
            leading=13.5,
            textColor=dark_color,
            spaceAfter=4
        )
        bullet_style = ParagraphStyle(
            'DocBullet',
            parent=body_style,
            leftIndent=14,
            firstLineIndent=-10,
            spaceAfter=2
        )
        callout_style = ParagraphStyle(
            'DocCalloutText',
            parent=body_style,
            fontSize=9,
            leading=13,
            textColor=dark_color
        )

        story = []

        # -------------------------------------------------------------
        # Cover Block / Header
        # -------------------------------------------------------------
        story.append(Paragraph(content.meta.title or "Test Engineering Specification", title_style))
        if content.meta.subtitle:
            story.append(Paragraph(content.meta.subtitle, sub_style))

        date_val = content.meta.date_str or datetime.now(timezone.utc).strftime("%B %d, %Y")
        meta_data = [
            [
                Paragraph(f"<b>Author:</b> {content.meta.author or 'EVAL AI'}", body_style),
                Paragraph(f"<b>Organization:</b> {content.meta.organization or 'Quality Engineering'}", body_style)
            ],
            [
                Paragraph(f"<b>Date:</b> {date_val}", body_style),
                Paragraph(f"<b>Classification:</b> {content.meta.classification or 'Confidential'}", body_style)
            ]
        ]
        meta_table = Table(meta_data, colWidths=[250, 250])
        meta_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, -1), light_bg),
            ('BOX', (0, 0), (-1, -1), 0.5, border_color),
            ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
            ('TOPPADDING', (0, 0), (-1, -1), 4),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 4),
            ('LEFTPADDING', (0, 0), (-1, -1), 8),
            ('RIGHTPADDING', (0, 0), (-1, -1), 8),
        ]))
        story.append(meta_table)
        story.append(Spacer(1, 10))

        # Executive Summary
        if content.executive_summary:
            summary_data = [[
                Paragraph(f"<b>EXECUTIVE SUMMARY</b><br/>{content.executive_summary}", callout_style)
            ]]
            summary_tbl = Table(summary_data, colWidths=[504])
            summary_tbl.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, -1), colors.HexColor("#F1F5F9")),
                ('BOX', (0, 0), (-1, -1), 1, secondary_color),
                ('TOPPADDING', (0, 0), (-1, -1), 8),
                ('BOTTOMPADDING', (0, 0), (-1, -1), 8),
                ('LEFTPADDING', (0, 0), (-1, -1), 10),
                ('RIGHTPADDING', (0, 0), (-1, -1), 10),
            ]))
            story.append(summary_tbl)
            story.append(Spacer(1, 10))

        # -------------------------------------------------------------
        # Sections
        # -------------------------------------------------------------
        for sec_idx, sec in enumerate(content.sections, start=1):
            story.append(Paragraph(f"{sec_idx}. {sec.heading}", h1_style))
            story.append(HRFlowable(width="100%", thickness=0.75, color=primary_color, spaceAfter=6))

            if sec.summary:
                story.append(Paragraph(f"<i>{sec.summary}</i>", body_style))

            for para in sec.paragraphs:
                story.append(Paragraph(para, body_style))

            for bullet in sec.bullet_points:
                story.append(Paragraph(f"&bull; {bullet}", bullet_style))

            # Callouts
            for callout in (sec.callouts or []):
                c_bg = colors.HexColor("#FEF3C7") if callout.type == "warning" else colors.HexColor("#EFF6FF")
                c_title = f"<b>[{callout.type.upper()}] {callout.title}</b><br/>" if callout.title else ""
                c_data = [[Paragraph(f"{c_title}{callout.content}", callout_style)]]
                c_tbl = Table(c_data, colWidths=[504])
                c_tbl.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, -1), c_bg),
                    ('BOX', (0, 0), (-1, -1), 0.5, border_color),
                    ('TOPPADDING', (0, 0), (-1, -1), 5),
                    ('BOTTOMPADDING', (0, 0), (-1, -1), 5),
                    ('LEFTPADDING', (0, 0), (-1, -1), 8),
                    ('RIGHTPADDING', (0, 0), (-1, -1), 8),
                ]))
                story.append(Spacer(1, 4))
                story.append(c_tbl)
                story.append(Spacer(1, 4))

            # Tables
            for tbl in (sec.tables or []):
                if tbl.caption:
                    story.append(Paragraph(f"<b>Table: {tbl.caption}</b>", body_style))

                table_rows = []
                if tbl.headers:
                    h_cells = [Paragraph(f"<b>{h}</b>", ParagraphStyle('TH', parent=body_style, textColor=colors.white, alignment=1)) for h in tbl.headers]
                    table_rows.append(h_cells)

                for r_data in tbl.rows:
                    r_cells = [Paragraph(str(val), body_style) for val in r_data]
                    table_rows.append(r_cells)

                if table_rows:
                    col_count = len(table_rows[0])
                    c_width = 504 / max(col_count, 1)
                    t_flow = Table(table_rows, colWidths=[c_width] * col_count)
                    t_style = [
                        ('BACKGROUND', (0, 0), (-1, 0), primary_color) if tbl.headers else ('BACKGROUND', (0, 0), (-1, -1), light_bg),
                        ('GRID', (0, 0), (-1, -1), 0.5, border_color),
                        ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
                        ('TOPPADDING', (0, 0), (-1, -1), 4),
                        ('BOTTOMPADDING', (0, 0), (-1, -1), 4),
                        ('LEFTPADDING', (0, 0), (-1, -1), 6),
                        ('RIGHTPADDING', (0, 0), (-1, -1), 6),
                    ]
                    # Zebra styling for data rows
                    start_r = 1 if tbl.headers else 0
                    for r_i in range(start_r, len(table_rows)):
                        if r_i % 2 == 1:
                            t_style.append(('BACKGROUND', (0, r_i), (-1, r_i), light_bg))
                    t_flow.setStyle(TableStyle(t_style))
                    story.append(Spacer(1, 4))
                    story.append(t_flow)
                    story.append(Spacer(1, 6))

        doc.build(story, canvasmaker=PageNumCanvas)
        stream.seek(0)
        return stream


# ===========================================================================
# 3. PPTX BUILDER (PowerPoint 16:9 Presentation Engine)
# ===========================================================================
class PptxBuilder:
    @classmethod
    def build(cls, content: DocumentContentModel, theme_name: str = "corporate_blue") -> io.BytesIO:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE

        theme = THEMES.get(theme_name, THEMES["corporate_blue"])
        p_rgb = hex_to_rgb(theme["primary_hex"])
        s_rgb = hex_to_rgb(theme["secondary_hex"])
        d_rgb = hex_to_rgb(theme["dark_hex"])
        l_rgb = hex_to_rgb(theme["light_bg_hex"])

        prs = Presentation()
        # Set 16:9 Widescreen dimensions (13.333 x 7.5 inches)
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank_slide_layout = prs.slide_layouts[6]

        slides = content.slides
        if not slides:
            # Fallback: construct slides from sections
            slides = [
                PptxSlide(
                    slide_number=1,
                    layout_type="title_slide",
                    title=content.meta.title or "Test Engineering Strategy & Report",
                    subtitle=content.meta.subtitle or "Automated Executive Briefing",
                    speaker_notes="Welcome to the executive test readiness review."
                )
            ]
            for idx, sec in enumerate(content.sections, start=2):
                slides.append(
                    PptxSlide(
                        slide_number=idx,
                        layout_type="card_grid",
                        title=sec.heading,
                        subtitle=sec.summary,
                        bullet_points=sec.bullet_points or sec.paragraphs[:3],
                        cards=[PptxSlideCard(title=f"Finding {i+1}", content=p[:120]) for i, p in enumerate(sec.paragraphs[:3])],
                        speaker_notes=f"Key discussion points for section {sec.heading}."
                    )
                )

        for slide_data in slides:
            slide = prs.slides.add_slide(blank_slide_layout)

            # Slide Background Accent
            bg_rect = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
            )
            bg_rect.fill.solid()
            bg_rect.fill.fore_color.rgb = RGBColor(*l_rgb)
            bg_rect.line.fill.background()

            # Top Accent Header Bar
            top_bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.12)
            )
            top_bar.fill.solid()
            top_bar.fill.fore_color.rgb = RGBColor(*p_rgb)
            top_bar.line.fill.background()

            # Speaker Notes
            if slide_data.speaker_notes:
                notes_slide = slide.notes_slide
                text_frame = notes_slide.notes_text_frame
                text_frame.text = slide_data.speaker_notes

            # ---------------------------------------------------------
            # Layout Type: Title Slide
            # ---------------------------------------------------------
            if slide_data.layout_type == "title_slide" or slide_data.slide_number == 1:
                # Big Corporate Banner
                banner = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.733), Inches(4.5)
                )
                banner.fill.solid()
                banner.fill.fore_color.rgb = RGBColor(*p_rgb)
                banner.line.fill.background()

                # Title Text
                t_box = slide.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(10.9), Inches(1.6))
                tf = t_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = slide_data.title or content.meta.title
                p.font.name = "Arial"
                p.font.size = Pt(36)
                p.font.bold = True
                p.font.color.rgb = RGBColor(255, 255, 255)

                # Subtitle
                if slide_data.subtitle or content.meta.subtitle:
                    s_box = slide.shapes.add_textbox(Inches(1.2), Inches(3.6), Inches(10.9), Inches(1.0))
                    stf = s_box.text_frame
                    stf.word_wrap = True
                    sp = stf.paragraphs[0]
                    sp.text = slide_data.subtitle or content.meta.subtitle or "Executive QA & Verification Deck"
                    sp.font.size = Pt(20)
                    sp.font.color.rgb = RGBColor(199, 210, 254)

                # Footer Meta in Banner
                m_box = slide.shapes.add_textbox(Inches(1.2), Inches(5.0), Inches(10.9), Inches(0.6))
                mtf = m_box.text_frame
                mp = mtf.paragraphs[0]
                date_val = content.meta.date_str or datetime.now(timezone.utc).strftime("%B %d, %Y")
                mp.text = f"{content.meta.author or 'EVAL AI'}  •  {content.meta.organization or 'Enterprise QA'}  •  {date_val}"
                mp.font.size = Pt(12)
                mp.font.color.rgb = RGBColor(224, 231, 255)
                continue

            # ---------------------------------------------------------
            # Standard Content Slide Header
            # ---------------------------------------------------------
            # Title
            t_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.45), Inches(11.7), Inches(0.8))
            tf = t_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = slide_data.title
            p.font.name = "Arial"
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = RGBColor(*p_rgb)

            # Subtitle
            if slide_data.subtitle:
                sp = tf.add_paragraph()
                sp.text = slide_data.subtitle
                sp.font.size = Pt(13)
                sp.font.color.rgb = RGBColor(100, 116, 139)

            # Running Footer
            foot_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.9), Inches(11.733), Inches(0.4))
            ftf = foot_box.text_frame
            fp = ftf.paragraphs[0]
            fp.text = f"EVAL AI  |  {content.meta.title[:45]}  |  Slide {slide_data.slide_number}"
            fp.font.size = Pt(9)
            fp.font.color.rgb = RGBColor(148, 163, 184)

            # ---------------------------------------------------------
            # Body: Cards Grid / Comparison / Bullet List
            # ---------------------------------------------------------
            cards = slide_data.cards or []
            if len(cards) > 0:
                card_count = min(len(cards), 3)
                card_width = (11.733 - (card_count - 1) * 0.4) / card_count
                card_height = 4.8

                for c_i, card in enumerate(cards[:card_count]):
                    c_left = 0.8 + c_i * (card_width + 0.4)
                    
                    # Card Background Shape
                    card_shape = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(c_left), Inches(1.6), Inches(card_width), Inches(card_height)
                    )
                    card_shape.fill.solid()
                    card_shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
                    card_shape.line.color.rgb = RGBColor(203, 213, 225)
                    card_shape.line.width = Pt(1)

                    # Card Top Highlight Strip
                    strip = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE, Inches(c_left), Inches(1.6), Inches(card_width), Inches(0.12)
                    )
                    strip.fill.solid()
                    strip.fill.fore_color.rgb = RGBColor(*s_rgb)
                    strip.line.fill.background()

                    # Card Text Frame
                    c_textbox = slide.shapes.add_textbox(
                        Inches(c_left + 0.25), Inches(1.85), Inches(card_width - 0.5), Inches(card_height - 0.5)
                    )
                    c_tf = c_textbox.text_frame
                    c_tf.word_wrap = True

                    # Card Title
                    cp = c_tf.paragraphs[0]
                    cp.text = card.title
                    cp.font.name = "Arial"
                    cp.font.bold = True
                    cp.font.size = Pt(14)
                    cp.font.color.rgb = RGBColor(*p_rgb)
                    cp.space_after = Pt(8)

                    # Card Content
                    body_p = c_tf.add_paragraph()
                    body_p.text = card.content
                    body_p.font.size = Pt(11)
                    body_p.font.color.rgb = RGBColor(*d_rgb)
            elif slide_data.bullet_points:
                # Bulleted Content Box
                b_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.0))
                b_tf = b_box.text_frame
                b_tf.word_wrap = True
                
                for bp_text in slide_data.bullet_points:
                    bp = b_tf.add_paragraph()
                    bp.text = f"•  {bp_text}"
                    bp.font.size = Pt(14)
                    bp.font.color.rgb = RGBColor(*d_rgb)
                    bp.space_after = Pt(10)

        stream = io.BytesIO()
        prs.save(stream)
        stream.seek(0)
        return stream


# ===========================================================================
# 4. DOC BUNDLE BUILDER (All Formats in Single ZIP)
# ===========================================================================
class DocBundleBuilder:
    @classmethod
    def build_zip(cls, content: DocumentContentModel, theme_name: str = "corporate_blue") -> io.BytesIO:
        docx_stream = DocxBuilder.build(content, theme_name)
        pdf_stream = PdfBuilder.build(content, theme_name)
        pptx_stream = PptxBuilder.build(content, theme_name)

        safe_title = (content.meta.title or "Test_Document").replace(" ", "_").replace("/", "_")
        zip_stream = io.BytesIO()
        with zipfile.ZipFile(zip_stream, "w", zipfile.ZIP_DEFLATED) as zf:
            zf.writestr(f"{safe_title}.docx", docx_stream.getvalue())
            zf.writestr(f"{safe_title}.pdf", pdf_stream.getvalue())
            zf.writestr(f"{safe_title}.pptx", pptx_stream.getvalue())

        zip_stream.seek(0)
        return zip_stream
