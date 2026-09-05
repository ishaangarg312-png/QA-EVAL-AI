import io
import csv
import logging
from typing import Optional

logger = logging.getLogger(__name__)

def extract_text_from_file(file_bytes: bytes, filename: str, max_chars: int = 50000) -> dict:
    """
    Extracts plaintext content from various file formats:
    - .pdf: using pypdf
    - .docx: using python-docx
    - .pptx: using python-pptx
    - .xlsx, .xls: using openpyxl
    - .csv, .tsv: standard csv parser
    - .txt, .md, .json: standard text decoder

    Returns a dict with:
      - filename
      - char_count
      - word_count
      - page_count / sheet_count (if applicable)
      - text: extracted plaintext
    """
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    extracted_text = ""
    extra_meta = {}

    try:
        if ext == "pdf":
            from pypdf import PdfReader
            stream = io.BytesIO(file_bytes)
            reader = PdfReader(stream)
            num_pages = len(reader.pages)
            pages_text = []
            for i, page in enumerate(reader.pages):
                txt = page.extract_text() or ""
                if txt.strip():
                    pages_text.append(f"--- [Page {i+1}] ---\n{txt}")
            extracted_text = "\n\n".join(pages_text)
            extra_meta["pages"] = num_pages

        elif ext in ("docx", "doc"):
            import docx
            stream = io.BytesIO(file_bytes)
            doc = docx.Document(stream)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            # Also extract tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        paragraphs.append(row_text)
            extracted_text = "\n\n".join(paragraphs)
            extra_meta["paragraphs"] = len(doc.paragraphs)

        elif ext in ("pptx", "ppt"):
            import pptx
            stream = io.BytesIO(file_bytes)
            prs = pptx.Presentation(stream)
            slide_texts = []
            for i, slide in enumerate(prs.slides):
                stexts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            if paragraph.text.strip():
                                stexts.append(paragraph.text.strip())
                if stexts:
                    slide_texts.append(f"--- [Slide {i+1}] ---\n" + "\n".join(stexts))
            extracted_text = "\n\n".join(slide_texts)
            extra_meta["slides"] = len(prs.slides)

        elif ext in ("xlsx", "xls"):
            import openpyxl
            stream = io.BytesIO(file_bytes)
            wb = openpyxl.load_workbook(stream, data_only=True)
            sheet_texts = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows_data = []
                for row in ws.iter_rows(values_only=True):
                    row_vals = [str(v).strip() for v in row if v is not None and str(v).strip()]
                    if row_vals:
                        rows_data.append(" | ".join(row_vals))
                if rows_data:
                    sheet_texts.append(f"--- [Sheet: {sheet_name}] ---\n" + "\n".join(rows_data[:500]))
            extracted_text = "\n\n".join(sheet_texts)
            extra_meta["sheets"] = len(wb.sheetnames)

        elif ext in ("csv", "tsv"):
            delimiter = "\t" if ext == "tsv" else ","
            try:
                decoded = file_bytes.decode("utf-8")
            except UnicodeDecodeError:
                decoded = file_bytes.decode("latin-1")
            reader = csv.reader(io.StringIO(decoded), delimiter=delimiter)
            rows = [" | ".join(cell.strip() for cell in row if cell.strip()) for row in reader]
            extracted_text = "\n".join(r for r in rows if r)

        else:
            # Fallback plain text / markdown / json
            try:
                extracted_text = file_bytes.decode("utf-8")
            except UnicodeDecodeError:
                extracted_text = file_bytes.decode("latin-1", errors="ignore")

    except Exception as e:
        logger.error(f"Error extracting text from {filename}: {str(e)}", exc_info=True)
        # Fallback to UTF-8 decoding if specific parser failed
        try:
            extracted_text = file_bytes.decode("utf-8", errors="ignore")
        except Exception:
            extracted_text = f"[Error reading {filename}: {str(e)}]"

    # Clean and truncate if too huge
    extracted_text = extracted_text.strip()
    original_length = len(extracted_text)
    if original_length > max_chars:
        extracted_text = extracted_text[:max_chars] + f"\n\n... [Truncated: Total {original_length} characters, limited to {max_chars}]"

    words = len(extracted_text.split())

    return {
        "filename": filename,
        "char_count": len(extracted_text),
        "word_count": words,
        "text": extracted_text,
        **extra_meta
    }
